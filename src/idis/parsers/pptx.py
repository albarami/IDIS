"""PPTX parser — deterministic text extraction with slide/shape locators.

Extracts text from PPTX files slide by slide, shape by shape,
producing SpanDraft objects with stable locators.

Requirements:
- Deterministic: same bytes in → same ordered spans out
- Fail-closed: malformed PPTX returns structured errors
- Stable locators: slides, shapes, paragraphs indexed in order
"""

from __future__ import annotations

import hashlib
import io
from typing import TYPE_CHECKING

from idis.parsers.base import (
    ParseError,
    ParseErrorCode,
    ParseLimits,
    ParseResult,
    SpanDraft,
)

if TYPE_CHECKING:
    pass

from pptx import Presentation
from pptx.exc import PackageNotFoundError


def _compute_content_hash(text: str) -> str:
    """Compute SHA-256 hash of text content."""
    return hashlib.sha256(text.encode("utf-8")).hexdigest()


def _normalize_text(text: str) -> str:
    """Normalize text for consistent output.

    - Convert \\r\\n and \\r to \\n
    - Strip trailing whitespace from each line
    - Strip leading/trailing whitespace from result
    """
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    lines = [line.rstrip() for line in text.split("\n")]
    return "\n".join(lines).strip()


def parse_pptx(
    data: bytes,
    limits: ParseLimits | None = None,
) -> ParseResult:
    """Parse PPTX bytes and extract text spans with slide/shape locators.

    Args:
        data: Raw PPTX file bytes.
        limits: Optional parsing limits (defaults to ParseLimits()).

    Returns:
        ParseResult with success=True and spans if extraction succeeded,
        or success=False with structured errors if parsing failed.

    Behavior:
        - Slides are processed in presentation order (0-indexed).
        - Shapes within each slide are processed in order.
        - Each text frame paragraph becomes a SpanDraft with locator
          {slide: s, shape: sh, paragraph: p}.
        - Table cells become SpanDraft with locator
          {slide: s, table: t, row: r, col: c}.
        - Empty paragraphs/cells are skipped.
        - Malformed PPTX files fail with CORRUPTED_FILE error.
    """
    if limits is None:
        limits = ParseLimits()

    if len(data) > limits.max_bytes:
        return ParseResult(
            doc_type="PPTX",
            success=False,
            errors=[
                ParseError(
                    code=ParseErrorCode.MAX_SIZE_EXCEEDED,
                    message=f"File size {len(data)} bytes exceeds limit {limits.max_bytes}",
                    details={"size": len(data), "limit": limits.max_bytes},
                )
            ],
        )

    try:
        prs = Presentation(io.BytesIO(data))
    except PackageNotFoundError as e:
        return ParseResult(
            doc_type="PPTX",
            success=False,
            errors=[
                ParseError(
                    code=ParseErrorCode.CORRUPTED_FILE,
                    message="Invalid PPTX file: not a valid Office Open XML package",
                    details={"error": str(e)},
                )
            ],
        )
    except Exception as e:
        return ParseResult(
            doc_type="PPTX",
            success=False,
            errors=[
                ParseError(
                    code=ParseErrorCode.CORRUPTED_FILE,
                    message="Failed to read PPTX file",
                    details={"error": str(e), "type": type(e).__name__},
                )
            ],
        )

    spans: list[SpanDraft] = []
    warnings: list[str] = []
    total_text_length = 0
    slide_count = len(prs.slides)
    shape_text_count = 0
    table_cell_count = 0

    for slide_idx, slide in enumerate(prs.slides):
        shape_idx = 0
        table_idx = 0

        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                for row_idx, row in enumerate(table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        text = _normalize_text(cell.text)
                        if not text:
                            continue

                        total_text_length += len(text)
                        table_cell_count += 1

                        spans.append(
                            SpanDraft(
                                span_type="CELL",
                                locator={
                                    "slide": slide_idx,
                                    "table": table_idx,
                                    "row": row_idx,
                                    "col": col_idx,
                                },
                                text_excerpt=text,
                                content_hash=_compute_content_hash(text),
                            )
                        )
                table_idx += 1

            elif shape.has_text_frame:
                text_frame = shape.text_frame
                for para_idx, paragraph in enumerate(text_frame.paragraphs):
                    text = _normalize_text(paragraph.text)
                    if not text:
                        continue

                    total_text_length += len(text)
                    shape_text_count += 1

                    spans.append(
                        SpanDraft(
                            span_type="PARAGRAPH",
                            locator={
                                "slide": slide_idx,
                                "shape": shape_idx,
                                "paragraph": para_idx,
                            },
                            text_excerpt=text,
                            content_hash=_compute_content_hash(text),
                        )
                    )
                shape_idx += 1

    if total_text_length == 0:
        return ParseResult(
            doc_type="PPTX",
            success=False,
            errors=[
                ParseError(
                    code=ParseErrorCode.NO_TEXT_EXTRACTED,
                    message="No extractable text found in PPTX",
                    details={"slide_count": slide_count},
                )
            ],
        )

    return ParseResult(
        doc_type="PPTX",
        success=True,
        spans=spans,
        metadata={
            "slide_count": slide_count,
            "shape_text_count": shape_text_count,
            "table_cell_count": table_cell_count,
            "span_count": len(spans),
            "total_text_length": total_text_length,
        },
        warnings=warnings,
    )
