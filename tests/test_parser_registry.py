"""Tests for parser registry — Phase 1.2.

Tests cover:
- Format detection by magic bytes (not extension/mime)
- Correct dispatch to PDF/XLSX/DOCX/PPTX parsers
- Fail-closed behavior for unknown formats
- Never raises exceptions on malformed input
"""

from __future__ import annotations

import io

import pytest
from docx import Document
from pptx import Presentation
from pptx.util import Inches

from idis.parsers.base import ParseErrorCode
from idis.parsers.registry import detect_format, parse_bytes

try:
    from openpyxl import Workbook
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas

    FIXTURES_AVAILABLE = True
except ImportError:
    FIXTURES_AVAILABLE = False


def create_test_pdf(text: str = "Test content") -> bytes:
    """Create a simple PDF for testing."""
    if not FIXTURES_AVAILABLE:
        pytest.skip("reportlab not installed")

    buffer = io.BytesIO()
    c = canvas.Canvas(buffer, pagesize=letter)
    c.drawString(72, 750, text)
    c.save()
    return buffer.getvalue()


def create_test_xlsx() -> bytes:
    """Create a simple XLSX for testing."""
    if not FIXTURES_AVAILABLE:
        pytest.skip("openpyxl not installed")

    wb = Workbook()
    ws = wb.active
    ws["A1"] = "Test"
    buffer = io.BytesIO()
    wb.save(buffer)
    return buffer.getvalue()


def create_test_docx(text: str = "Test content") -> bytes:
    """Create a simple DOCX for testing."""
    doc = Document()
    doc.add_paragraph(text)
    buffer = io.BytesIO()
    doc.save(buffer)
    return buffer.getvalue()


def create_test_pptx(text: str = "Test content") -> bytes:
    """Create a simple PPTX for testing."""
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    txbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(0.5))
    txbox.text_frame.text = text
    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


class TestFormatDetection:
    """Test format detection by magic bytes."""

    def test_detect_pdf_by_magic(self) -> None:
        """PDF detected by %PDF- magic bytes."""
        pdf_bytes = create_test_pdf()

        detected = detect_format(pdf_bytes)

        assert detected == "PDF"

    def test_detect_pdf_ignores_extension(self) -> None:
        """PDF detection based on content, not filename extension."""
        pdf_bytes = create_test_pdf()

        result = parse_bytes(pdf_bytes, filename="document.xlsx")

        assert result.doc_type == "PDF"
        assert result.success is True

    def test_detect_xlsx_by_zip_structure(self) -> None:
        """XLSX detected by ZIP with xl/workbook.xml."""
        xlsx_bytes = create_test_xlsx()

        detected = detect_format(xlsx_bytes)

        assert detected == "XLSX"

    def test_detect_xlsx_ignores_extension(self) -> None:
        """XLSX detection based on content, not filename extension."""
        xlsx_bytes = create_test_xlsx()

        result = parse_bytes(xlsx_bytes, filename="document.pdf")

        assert result.doc_type == "XLSX"
        assert result.success is True

    def test_unknown_format_returns_none(self) -> None:
        """Unknown format returns None from detect_format."""
        random_bytes = b"This is just random text data"

        detected = detect_format(random_bytes)

        assert detected is None

    def test_zip_without_workbook_not_xlsx(self) -> None:
        """ZIP file without xl/workbook.xml is not detected as XLSX."""
        import zipfile

        buffer = io.BytesIO()
        with zipfile.ZipFile(buffer, "w") as zf:
            zf.writestr("data.txt", "Hello world")

        zip_bytes = buffer.getvalue()
        detected = detect_format(zip_bytes)

        assert detected is None

    def test_detect_docx_by_zip_structure(self) -> None:
        """DOCX detected by ZIP with word/document.xml."""
        docx_bytes = create_test_docx()

        detected = detect_format(docx_bytes)

        assert detected == "DOCX"

    def test_detect_pptx_by_zip_structure(self) -> None:
        """PPTX detected by ZIP with ppt/presentation.xml."""
        pptx_bytes = create_test_pptx()

        detected = detect_format(pptx_bytes)

        assert detected == "PPTX"


class TestParseBytes:
    """Test unified parse_bytes entrypoint."""

    def test_parse_pdf(self) -> None:
        """parse_bytes correctly dispatches to PDF parser."""
        pdf_bytes = create_test_pdf("Test content")

        result = parse_bytes(pdf_bytes)

        assert result.success is True
        assert result.doc_type == "PDF"
        assert len(result.spans) > 0

    def test_parse_xlsx(self) -> None:
        """parse_bytes correctly dispatches to XLSX parser."""
        xlsx_bytes = create_test_xlsx()

        result = parse_bytes(xlsx_bytes)

        assert result.success is True
        assert result.doc_type == "XLSX"
        assert len(result.spans) > 0

    def test_parse_docx(self) -> None:
        """parse_bytes correctly dispatches to DOCX parser."""
        docx_bytes = create_test_docx("Test content")

        result = parse_bytes(docx_bytes)

        assert result.success is True
        assert result.doc_type == "DOCX"
        assert len(result.spans) > 0

    def test_parse_pptx(self) -> None:
        """parse_bytes correctly dispatches to PPTX parser."""
        pptx_bytes = create_test_pptx("Test content")

        result = parse_bytes(pptx_bytes)

        assert result.success is True
        assert result.doc_type == "PPTX"
        assert len(result.spans) > 0

    def test_unsupported_format_error(self) -> None:
        """Unsupported format returns structured error."""
        random_bytes = b"This is not a supported format"

        result = parse_bytes(random_bytes)

        assert result.success is False
        assert len(result.errors) > 0
        error_codes = [e.code for e in result.errors]
        assert ParseErrorCode.UNSUPPORTED_FORMAT in error_codes

    def test_empty_file_error(self) -> None:
        """Empty file returns structured error with UNKNOWN doc_type."""
        result = parse_bytes(b"")

        assert result.success is False
        assert result.doc_type == "UNKNOWN"
        assert len(result.errors) > 0
        assert result.errors[0].code == ParseErrorCode.UNSUPPORTED_FORMAT

    def test_unsupported_format_has_unknown_doc_type(self) -> None:
        """Unsupported format returns UNKNOWN doc_type."""
        result = parse_bytes(b"random data that is not a valid format")

        assert result.success is False
        assert result.doc_type == "UNKNOWN"


class TestRegistryNeverRaises:
    """Test that registry never raises exceptions on any input."""

    def test_random_bytes_no_exception(self) -> None:
        """Random bytes don't raise exception."""
        result = parse_bytes(b"\x00\x01\x02\xff\xfe\xfd" * 100)

        assert result.success is False

    def test_partial_pdf_header_no_exception(self) -> None:
        """Partial PDF header doesn't raise exception."""
        result = parse_bytes(b"%PDF-1.4\n" + b"\xff" * 50)

        assert result.success is False

    def test_partial_zip_header_no_exception(self) -> None:
        """Partial ZIP header doesn't raise exception."""
        result = parse_bytes(b"PK\x03\x04" + b"\xff" * 50)

        assert result.success is False

    def test_very_short_input_no_exception(self) -> None:
        """Very short input doesn't raise exception."""
        result = parse_bytes(b"PK")

        assert result.success is False

    def test_null_bytes_no_exception(self) -> None:
        """Null bytes don't raise exception."""
        result = parse_bytes(b"\x00" * 1000)

        assert result.success is False


class TestRegistryErrorDetails:
    """Test error details in registry responses."""

    def test_unsupported_format_includes_header_bytes(self) -> None:
        """Unsupported format error includes header bytes for debugging."""
        test_bytes = b"UNKNOWN_FORMAT_HEADER_12345"

        result = parse_bytes(test_bytes, filename="test.bin")

        assert result.success is False
        error = result.errors[0]
        assert error.code == ParseErrorCode.UNSUPPORTED_FORMAT
        assert "header_bytes" in error.details

    def test_unsupported_format_includes_filename(self) -> None:
        """Unsupported format error includes filename if provided."""
        test_bytes = b"not a valid format"

        result = parse_bytes(test_bytes, filename="mystery.dat")

        assert result.success is False
        error = result.errors[0]
        assert error.details.get("filename") == "mystery.dat"

    def test_unsupported_format_includes_mime_type(self) -> None:
        """Unsupported format error includes MIME type if provided."""
        test_bytes = b"not a valid format"

        result = parse_bytes(test_bytes, mime_type="application/octet-stream")

        assert result.success is False
        error = result.errors[0]
        assert error.details.get("mime_type") == "application/octet-stream"


class TestRegistryIntegration:
    """Integration tests for registry with real parsers."""

    def test_pdf_spans_have_valid_locators(self) -> None:
        """PDF spans from registry have valid page/line locators."""
        pdf_bytes = create_test_pdf("Integration test")

        result = parse_bytes(pdf_bytes)

        assert result.success is True
        for span in result.spans:
            assert span.span_type == "PAGE_TEXT"
            assert "page" in span.locator
            assert "line" in span.locator

    def test_xlsx_spans_have_valid_locators(self) -> None:
        """XLSX spans from registry have valid sheet/cell locators."""
        xlsx_bytes = create_test_xlsx()

        result = parse_bytes(xlsx_bytes)

        assert result.success is True
        for span in result.spans:
            assert span.span_type == "CELL"
            assert "sheet" in span.locator
            assert "cell" in span.locator
            assert "row" in span.locator
            assert "col" in span.locator

    def test_docx_spans_have_valid_locators(self) -> None:
        """DOCX spans from registry have valid paragraph locators."""
        docx_bytes = create_test_docx("Integration test")

        result = parse_bytes(docx_bytes)

        assert result.success is True
        for span in result.spans:
            assert span.span_type == "PARAGRAPH"
            assert "paragraph" in span.locator

    def test_pptx_spans_have_valid_locators(self) -> None:
        """PPTX spans from registry have valid slide/shape locators."""
        pptx_bytes = create_test_pptx("Integration test")

        result = parse_bytes(pptx_bytes)

        assert result.success is True
        for span in result.spans:
            assert "slide" in span.locator
            if span.span_type == "PARAGRAPH":
                assert "shape" in span.locator
                assert "paragraph" in span.locator
