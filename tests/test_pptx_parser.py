"""Tests for PPTX parser — deterministic extraction with stable locators.

Tests verify:
- Successful parsing with correct locators
- Deterministic output (same input → same spans)
- Fail-closed behavior for corrupted files
- Table extraction with slide/table/row/col locators
"""

from __future__ import annotations

import io

from pptx import Presentation
from pptx.util import Inches

from idis.parsers.base import ParseErrorCode, ParseLimits
from idis.parsers.pptx import parse_pptx


def create_test_pptx(
    slides: list[list[str]],
    tables: list[tuple[int, list[list[str]]]] | None = None,
) -> bytes:
    """Create a PPTX file in memory with specified slide content.

    Args:
        slides: List of slides, each containing list of text box contents.
        tables: Optional list of (slide_index, table_data) tuples.

    Returns:
        PPTX file as bytes.
    """
    prs = Presentation()

    for slide_texts in slides:
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        left = Inches(1)
        top = Inches(1)
        for i, text in enumerate(slide_texts):
            txbox = slide.shapes.add_textbox(left, top + Inches(i * 0.5), Inches(8), Inches(0.5))
            tf = txbox.text_frame
            tf.text = text

    if tables:
        for slide_idx, table_data in tables:
            if slide_idx >= len(prs.slides):
                continue
            slide = prs.slides[slide_idx]
            num_rows = len(table_data)
            num_cols = max(len(row) for row in table_data) if table_data else 0
            if num_cols == 0:
                continue

            table = slide.shapes.add_table(
                num_rows, num_cols, Inches(1), Inches(4), Inches(6), Inches(2)
            ).table
            for row_idx, row_data in enumerate(table_data):
                for col_idx, cell_text in enumerate(row_data):
                    table.cell(row_idx, col_idx).text = cell_text

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


class TestPPTXParserSuccess:
    """Test successful PPTX parsing."""

    def test_parse_simple_pptx(self) -> None:
        """Parse PPTX with text boxes returns success."""
        slides = [["Title slide"], ["Content slide"]]
        pptx_bytes = create_test_pptx(slides)

        result = parse_pptx(pptx_bytes)

        assert result.success is True
        assert result.doc_type == "PPTX"
        assert len(result.errors) == 0
        assert len(result.spans) >= 2

    def test_parse_pptx_with_tables(self) -> None:
        """Parse PPTX with tables extracts cells."""
        slides = [["Header"]]
        tables = [(0, [["A1", "B1"], ["A2", "B2"]])]
        pptx_bytes = create_test_pptx(slides, tables)

        result = parse_pptx(pptx_bytes)

        assert result.success is True
        cell_spans = [s for s in result.spans if s.span_type == "CELL"]
        assert len(cell_spans) == 4

    def test_metadata_populated(self) -> None:
        """Metadata contains slide and span counts."""
        slides = [["Slide 1"], ["Slide 2"], ["Slide 3"]]
        pptx_bytes = create_test_pptx(slides)

        result = parse_pptx(pptx_bytes)

        assert result.metadata["slide_count"] == 3
        assert result.metadata["span_count"] >= 3
        assert result.metadata["total_text_length"] > 0


class TestPPTXParserLocators:
    """Test locator correctness for PPTX spans."""

    def test_text_locators_include_slide_shape_paragraph(self) -> None:
        """Text span locators include slide, shape, paragraph."""
        slides = [["First text", "Second text"]]
        pptx_bytes = create_test_pptx(slides)

        result = parse_pptx(pptx_bytes)

        assert result.success is True
        para_spans = [s for s in result.spans if s.span_type == "PARAGRAPH"]
        assert len(para_spans) >= 2

        for span in para_spans:
            assert "slide" in span.locator
            assert "shape" in span.locator
            assert "paragraph" in span.locator

    def test_table_cell_locators(self) -> None:
        """Table cell locators include slide, table, row, col."""
        slides = [[]]
        tables = [(0, [["R0C0", "R0C1"], ["R1C0", "R1C1"]])]
        pptx_bytes = create_test_pptx(slides, tables)

        result = parse_pptx(pptx_bytes)

        assert result.success is True
        cell_spans = [s for s in result.spans if s.span_type == "CELL"]
        assert len(cell_spans) == 4

        for span in cell_spans:
            assert "slide" in span.locator
            assert "table" in span.locator
            assert "row" in span.locator
            assert "col" in span.locator

    def test_slide_indices_zero_based(self) -> None:
        """Slide indices are 0-based."""
        slides = [["Slide 0"], ["Slide 1"]]
        pptx_bytes = create_test_pptx(slides)

        result = parse_pptx(pptx_bytes)

        assert result.success is True
        slide_indices = {s.locator["slide"] for s in result.spans}
        assert 0 in slide_indices
        assert 1 in slide_indices

    def test_content_hash_populated(self) -> None:
        """Each span has a content hash."""
        slides = [["Test content"]]
        pptx_bytes = create_test_pptx(slides)

        result = parse_pptx(pptx_bytes)

        assert result.success is True
        for span in result.spans:
            assert span.content_hash is not None
            assert len(span.content_hash) == 64  # SHA-256 hex


class TestPPTXParserDeterminism:
    """Test deterministic output."""

    def test_same_input_same_output(self) -> None:
        """Parsing same bytes twice produces identical spans."""
        slides = [["Revenue: $10M"], ["Growth: 150%"]]
        pptx_bytes = create_test_pptx(slides)

        result1 = parse_pptx(pptx_bytes)
        result2 = parse_pptx(pptx_bytes)

        assert result1.success is True
        assert result2.success is True
        assert len(result1.spans) == len(result2.spans)

        for span1, span2 in zip(result1.spans, result2.spans, strict=True):
            assert span1.locator == span2.locator
            assert span1.text_excerpt == span2.text_excerpt
            assert span1.content_hash == span2.content_hash
            assert span1.span_type == span2.span_type


class TestPPTXParserFailClosed:
    """Test fail-closed behavior for invalid inputs."""

    def test_corrupted_bytes(self) -> None:
        """Corrupted data returns error, no exception."""
        corrupted = b"PK\x03\x04not_a_real_pptx_file"

        result = parse_pptx(corrupted)

        assert result.success is False
        assert result.doc_type == "PPTX"
        assert len(result.errors) > 0
        assert result.errors[0].code == ParseErrorCode.CORRUPTED_FILE

    def test_random_bytes(self) -> None:
        """Random bytes return corrupted file error."""
        random_data = b"\x00\x01\x02\x03\x04\x05\x06\x07\x08\x09"

        result = parse_pptx(random_data)

        assert result.success is False
        assert result.errors[0].code == ParseErrorCode.CORRUPTED_FILE

    def test_empty_pptx_no_text(self) -> None:
        """PPTX with no text returns NO_TEXT_EXTRACTED."""
        pptx_bytes = create_test_pptx([[]])

        result = parse_pptx(pptx_bytes)

        assert result.success is False
        assert result.errors[0].code == ParseErrorCode.NO_TEXT_EXTRACTED


class TestPPTXParserLimits:
    """Test limit enforcement."""

    def test_max_size_exceeded(self) -> None:
        """File exceeding max_bytes returns error."""
        slides = [["Test slide"]]
        pptx_bytes = create_test_pptx(slides)

        limits = ParseLimits(max_bytes=100)  # Very small limit
        result = parse_pptx(pptx_bytes, limits=limits)

        assert result.success is False
        assert result.errors[0].code == ParseErrorCode.MAX_SIZE_EXCEEDED


class TestPPTXParserEdgeCases:
    """Test edge cases."""

    def test_empty_text_boxes_skipped(self) -> None:
        """Empty text boxes are not included in spans."""
        slides = [["First", "", "Third"]]
        pptx_bytes = create_test_pptx(slides)

        result = parse_pptx(pptx_bytes)

        assert result.success is True
        texts = [s.text_excerpt for s in result.spans]
        assert "First" in texts
        assert "Third" in texts
        assert "" not in texts

    def test_unicode_content(self) -> None:
        """Unicode content is preserved."""
        slides = [["Hello 世界", "مرحبا بالعالم"]]
        pptx_bytes = create_test_pptx(slides)

        result = parse_pptx(pptx_bytes)

        assert result.success is True
        texts = [s.text_excerpt for s in result.spans]
        assert "Hello 世界" in texts
        assert "مرحبا بالعالم" in texts

    def test_multiple_slides(self) -> None:
        """Multiple slides are all extracted."""
        slides = [["Slide 1 content"], ["Slide 2 content"], ["Slide 3 content"]]
        pptx_bytes = create_test_pptx(slides)

        result = parse_pptx(pptx_bytes)

        assert result.success is True
        slide_indices = {s.locator["slide"] for s in result.spans}
        assert len(slide_indices) == 3
